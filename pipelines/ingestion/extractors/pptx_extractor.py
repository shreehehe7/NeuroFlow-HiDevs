import os
from pptx import Presentation
from PIL import Image
import io
from typing import List
from google import genai
from . import ExtractedPage

def extract_pptx(file_path: str) -> List[ExtractedPage]:
    pages: List[ExtractedPage] = []
    prs = Presentation(file_path)
    
    api_key = os.getenv("GEMINI_API_KEY")
    client = genai.Client(api_key=api_key) if api_key else None
    
    for i, slide in enumerate(prs.slides):
        page_number = i + 1
        slide_text = []
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
                
        # Extract speaker notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            
        combined_text = "\n".join(slide_text)
        if notes:
            combined_text += f"\n\nSpeaker Notes:\n{notes}"
            
        # Check for images and use Vision LLM
        vision_description = ""
        for shape in slide.shapes:
            if shape.shape_type == 13: # 13 is msoPicture
                try:
                    image_stream = io.BytesIO(shape.image.blob)
                    img = Image.open(image_stream)
                    # Resize to max 1024px
                    if max(img.size) > 1024:
                        img.thumbnail((1024, 1024), Image.Resampling.LANCZOS)
                        
                    if client:
                        prompt = "Describe this diagram or image from a presentation slide in detail."
                        response = client.models.generate_content(
                            model='gemini-1.5-flash',
                            contents=[img, prompt]
                        )
                        vision_description += f"\n\nImage Description: {response.text}"
                except Exception as e:
                    vision_description += f"\n\n[Failed to extract image description: {str(e)}]"
                    
        final_content = combined_text + vision_description
        
        pages.append(ExtractedPage(
            page_number=page_number,
            content=final_content.strip(),
            content_type="text",
            metadata={"has_notes": bool(notes), "has_vision": bool(vision_description)}
        ))
        
    return pages
